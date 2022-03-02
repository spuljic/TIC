#! /usr/bin/env python
'''
TODO:
max cycle dynamic for cfx 20211116
set a minimum number of cicles to accept analysis


analyze single run or historical run.
you have to run a historical folder to then move on and run analysis_historic.py
python3 anal_pcr.py -i input file
python3 anal_pcr.py -f input folder

#folder strcuture should be part of this
HISTORICAL_ANALYSIS_DATA
    Data_ABI7500
        CENTER_X
            RUN_XY
                data file
    Data_CFX96
        CENTER_Y
            RUN_YZ
                data file
you can run all the tree or part of it

mandatory folder names
ABI principal folder must HAVE "ABI7500" pattern on its name
CFX principal folder must HAVE "CFX96" pattern on its name

cfx machine yield several files with data we are only interested in
    *Quanti.*tion Amplification Results.xlsx --> is an xlsx file
    ex:
    admin_2021-03-31 13-04-33_DX100054 -  Quantitation Amplification Results.xlsx
for abi machines we are interested in 
    *data.xls --> is an xls file
    ex:
    TIC15022021_data.xls

output:
for each file analyzed several files are generated
output.date.sigmoid.pptx: is a summary in power point to share with people involved in the project
                          it includes:
                          -scatter plots
                          -sigmoid plots
                          -fluorescence vs amplimix
                          it should include genotyping results

output.date.raw.{pkl, xlsx}: raw dataframe in excel and pickle format
columns: mode	center	runcode	Well	Cycle	Rn	ΔRn	Sample Name	Reporter
	Amplimix	Factor

output.date.final.{pkl, xlsx}: final dataframe in excel and pickle format
columns: mode	center	runcode	Well	Sample Name	Amplimix	ΔRn_FAM	
    FAM_uas	FAM_supera_uas	ΔRn_HEX	HEX_uas	HEX_supera_uas	Ratio_FAM/HEX
    NoCall_FAM/HEX	ΔRn_TEX	TEX_uas	TEX_supera_uas	ΔRn_ATTO	ATTO_uas
    ATTO_supera_uas	Ratio_TEX/ATTO	NoCall_TEX/ATTO	Empty

output: is the root input file name with amendments to avoid problems 
(get rid of white spaces, non compatible chars, lower_case)
date: yyyymmdd for the analysis day
additionally each run ends up with a concatenation of pkl
final_total.pkl --> this file is the input for the analysis_historic.pkl script
raw_total.pkl

'''
import argparse
import sys
import os
import re
import time
import unicodedata
import sqlite3
import math
import pprint

from shutil import rmtree, move
from collections import defaultdict

import pandas as pd 
import numpy as np

from matplotlib import pyplot
import seaborn as sns
from matplotlib.lines import Line2D  # for legend handle
from pptx import Presentation
from pptx.util import Cm, Pt

def main(args):
    '''do everything
    '''
    db_folder = './DB/'
    png_folder = 'PNG_TMP/'
    if not os.path.exists(png_folder):
        os.makedirs(png_folder)
    else:
        rmtree(png_folder)
        os.makedirs(png_folder)
    prefix = 'ΔRn'
    #prefix = 'Rn'
    patterns = {
        'cfx_pat': r'.*CFX96/',
        'abi_pat': r'.*ABI7500/',
        'cfx_file_path': r' -  Quanti.*tion Amplification Results',
        'abi_file_path': r'_data'
    }
    '''
    (ADD_HISTORIC, LOAD_HISTORIC, db_folder, png_folder, out_file_root,
     timestamp, historic_file, backup_file) = setup()
    historic_list = []
    '''
    out_file_root = db_folder + 'historic_ratios'
    historic_file = '{}.pkl'.format(out_file_root)
    LOAD_HISTORIC = False
    df_historic = load_historic(LOAD_HISTORIC, historic_file, prefix)
    #load db folder
    list_file_in = []
    if not os.path.exists(db_folder):
        print('Error: database folder: {} not found.'.format(db_folder))
        sys.exit()
    if args.infile:
        center = '.'#in order to manage historic data we should pass center(free text or list)
        list_file_in = load_input_file(args.infile, patterns, center)
        if len(list_file_in) < 1:
            print('input file not valid')
            sys.exit()
        final_list, raw_list = analyze_file(db_folder,  list_file_in, prefix, df_historic, png_folder)
    elif args.folder:
        #manage folder
        #check if it is a folder
        if not os.path.isdir(args.folder):
            print('please provide a valid folder as input')
            sys.exit()
        list_file_in = load_folder(args.folder, patterns)
        #check if at least one file has been load
        if len(list_file_in) < 1:#no files load
            print("no files load from folder {}, nothing to process".format(args.folder))
            sys.exit()
        final_list, raw_list = analyze_file(db_folder,  list_file_in, prefix, df_historic, png_folder)
        df_final_total = pd.concat(final_list, axis=0, ignore_index=True)
        df_raw_total = pd.concat(raw_list, axis=0, ignore_index=True)
        df_final_total.to_pickle('final_total.pkl')
        df_raw_total.to_pickle('raw_total.pkl')
    else:
        print('Define either file or folder')
        sys.exit()
    '''
    if ADD_HISTORIC:
        print(historic_list)
        df_historic_new = pd.concat(historic_list, axis=0, ignore_index=True)
        if os.path.isfile(historic_file):#backup file
            move(historic_file, backup_file)
            print('Info: backuping historic file to {}'.format(backup_file))
        df_historic_new.to_pickle(historic_file)
        print('Info: historic saved to file to {}'.format(historic_file))
    '''
###############################################################################
def analyze_file(db_folder, list_file_in, prefix, df_historic, png_folder):
    ''' from a list of files extract data
    '''
    final_list = list()
    raw_list = list()
    ppt_file = True
    timestr = time.strftime("%Y%m%d") 
    for sample_dict in list_file_in:
        #load metadatae
        mode = sample_dict['mode']
        center = sample_dict['center']
        runcode = sample_dict['runcode']
        filein  = sample_dict['filepath']
        #folder = sample_dict['folder']
        #load proper database with ratio and uas
        db_file = load_db_plugin(mode, db_folder)   
        df_ratio = load_ratio_from_db(db_file, mode)
        df_uas = load_uas_from_db(db_file, mode)
        if mode == 'ABI':#load abi data
            df_final, df_raw = process_abi_file(prefix, filein, df_uas,
                                                db_file, mode)
        elif mode == 'CFX':#load cfx data
            df_final, df_raw = process_cfx_file(prefix, filein, df_uas, mode)
        else:
            print("mode: {} doesn't exist".format(mode))
            sys.exit()
        #adding columns
        df_final.insert(loc = 0, column = 'mode', value = mode)
        df_final.insert(loc = 1, column = 'center', value = center)
        df_final.insert(loc = 2, column = 'runcode', value = runcode)
        df_raw.insert(loc = 0, column = 'mode', value = mode)
        df_raw.insert(loc = 1, column = 'center', value = center)
        df_raw.insert(loc = 2, column = 'runcode', value = runcode)
        #todo
        #do passing uas analysis per sample, row and column
        #identify ctrl samples by position (actual) but also identify by profile
        #this is:
        #ctrl- no uas for most of the wells/flour
        #filter out low quality samples before analysis (ie: outliers of some kind, non passing uas samples) 
        #plot all samples vs good quality samples
        #add final genotype by ratios (actual method)
        #explore genotyping by clustering:
        # (test several clustering methods):
        # it seems important to know the number of clusters which is a probability linked to allele freq
        # and number of samples tested 
        if(ppt_file):#yield a ppt file
            pptx_file_name = '{}.{}.sigmoid.pptx'.format(runcode, timestr)
            prs = Presentation()
            prs.slide_width = Cm(33.867)
            prs.slide_height = Cm(19.05)
            plot_ratios(prs, prefix, df_final, df_ratio, df_uas, df_historic, runcode, mode, png_folder)
            plot_fluor(prs, prefix, df_final, df_ratio, df_uas, center, mode, png_folder)
            plot_data_plate(prs, runcode, timestr, df_raw, df_uas, mode, png_folder)
            prs.save(pptx_file_name)
        #final oputput file
        df_final.to_pickle('{}.{}.final.pkl'.format(runcode, timestr))
        df_raw.to_pickle('{}.{}.raw.pkl'.format(runcode, timestr))
        df_final.to_excel('{}.{}.final.xlsx'.format(runcode, timestr))
        df_raw.to_excel('{}.{}.raw.xlsx'.format(runcode, timestr))
        final_list.append(df_final)
        raw_list.append(df_raw)
    return final_list, raw_list
###############################################################################
def load_folder(folder, patterns):
    #walk through files and get the metadata to load data later on
    #end up with a list of dictionaries with metadata
    list_of_files = list()
    for (dirpath, dirnames, filenames) in os.walk(folder):
        print(dirpath)
        if len(filenames) > 0:
            for filename in filenames:
                f_root, f_ext = os.path.splitext(filename)
                if f_ext == '.xls':
                    mode = 'ABI'
                    #convert file name
                    run = slugify(re.sub(patterns['abi_file_path'], '', f_root))
                    center = slugify(re.sub(r'/.*', '', re.sub(patterns['abi_pat'], '', dirpath)))
                elif f_ext == '.xlsx' and re.search(patterns['cfx_file_path'], f_root):
                    mode = 'CFX'
                    run = slugify(re.sub(patterns['cfx_file_path'], '', f_root))
                    center = slugify(re.sub(r'/.*', '', re.sub(patterns['cfx_pat'], '', dirpath)))
                else:
                    continue
                f_path = os.sep.join([dirpath, filename])
                list_of_files.append({'mode': mode, 'runcode': run,
                                      'center': center, 'filename': filename,
                                      'filepath': f_path, 'folder': dirpath})
    #pprint.pprint(list_of_files)
    return list_of_files
###############################################################################
def load_input_file(infile, patterns, center):
    #load a single file input
    list_of_files = []
    filename = os.path.basename(infile.name)
    dirpath = os.path.dirname(infile.name)
    if not dirpath:
        dirpath = './'
    f_root, f_ext = os.path.splitext(filename)
    if f_ext == '.xls':
        mode = 'ABI'
        run = slugify(re.sub(patterns['abi_file_path'], '', f_root))
    elif f_ext == '.xlsx':
        mode = 'CFX'
        run = slugify(re.sub(patterns['cfx_file_path'], '', f_root))
    else:
        print('extension not accepted')
        sys.exit()
    f_path =os.sep.join([dirpath, filename])
    list_of_files.append({'mode': mode, 'runcode': run, 'center': center, 
                          'filename': filename, 'filepath': f_path,
                          'folder': dirpath})
    return list_of_files
###############################################################################
def load_historic(LOAD_HISTORIC, historic_file, prefix):
    #this function is deprecated
    #it would like to manage historic data and plot it  
    if LOAD_HISTORIC:
        #ha de funcionar pels 2 modes
        if os.path.isfile(historic_file):#backup file
            print('Info: loading historic file:{}'.format(historic_file))
            df_historic = pd.read_pickle(historic_file)
            historic_list.append(df_historic)
            #df = df.drop(columns=['index'])
            #drop index col
        else:
            print('Error: historic file:{} not found'.format(historic_file))
            LOAD_HISTORIC = False
            #sys.exit()
    else:
        #create empty df 
        df_historic = pd.DataFrame(columns=['Well', 'Sample Name',
            'Amplimix', prefix + '_FAM', 'FAM_uas',
            'FAM_supera_uas', prefix + '_HEX', 'HEX_uas', 'HEX_supera_uas',
            'Ratio_FAM/HEX', 'NoCall_FAM/HEX', prefix + '_TEX', 'TEX_uas',
            'TEX_supera_uas', prefix + '_ATTO', 'ATTO_uas', 'ATTO_supera_uas',
            'Ratio_TEX/ATTO', 'NoCall_TEX/ATTO', 'Empty'])
    return df_historic
###############################################################################
'''
def setup():
    #set operational folders and files
    ADD_HISTORIC = False
    LOAD_HISTORIC = False
    db_folder = './DB/'
    png_folder = 'PNG_TMP/'
    if not os.path.exists(png_folder):
        os.makedirs(png_folder)
    else:
        rmtree(png_folder)
        os.makedirs(png_folder)
    out_file_root = db_folder + 'historic_ratios'
    timestamp = time.strftime("%Y%m%d-%H%M%S") 
    historic_file = '{}.pkl'.format(out_file_root)
    backup_file = '{}.{}.pkl'.format(out_file_root, timestamp)
    return(ADD_HISTORIC, LOAD_HISTORIC, db_folder, png_folder, 
           out_file_root, timestamp, historic_file, backup_file)
'''
###############################################################################
def process_cfx_file(prefix, filein, df_uas, mode):
    '''process all cfx file
    '''
    #extract data from cfx file *Quantification Amplification Results.xlsx
    df_work, df_raw = extract_cfx(prefix, filein)
    df_final = manage_work(prefix, df_work, df_uas, mode)
    return df_final, df_raw
###############################################################################
def process_abi_file(prefix, filein, df_uas, db_file, mode):
    '''process all abi file
    '''
    #extract data from abi
    df_work, df_raw = extract_abi(prefix, filein, db_file)
    df_final = manage_work(prefix, df_work, df_uas, mode)
    return df_final, df_raw
###############################################################################
def load_db_plugin(mode, db_folder):
    '''load thresholds from plugin sqlite
    in order to avoid problems with data source and manage older data 
    '''
    #en funció d'abi o de cfx
    db_new_file = ''
    db_new = ''
    if mode == 'ABI':
        #db_old_file = 'TIC_03_ABI.old.sqlite' 
        #db_old = db_folder + db_old_file
        db_new_file = 'TIC_04_ABI.sqlite'
        db_new = db_folder + db_new_file
        #db for previous version
        #db_new_file = 'TIC_03_ABI.old.sqlite' 
        #db_new = db_folder + db_new_file
    elif mode == 'CFX':
        db_new_file = 'TIC_03_CFX'
        db_new = db_folder + db_new_file
    for dbfile in ([db_new]):
        if not os.path.isfile(dbfile):
            print('Error: database file: {} not found.'.format(dbfile))
            sys.exit()
    return db_new
###############################################################################
def plot_ratios(prs, prefix, df_final, df_ratio, df_uas, df_historic, root, mode, folder):
    '''add slides of scatter plots
    '''
    amp2 = False
    for amplimix in df_final['Amplimix'].unique():
        if amplimix == 2:
            amp2 = True
        #fl1 = 'FAM'
        #fl2 = 'HEX'
        fl2 = 'FAM'
        fl1 = 'HEX'
        fig2 = plot_figure_ratio(prefix, df_final, df_ratio, df_uas, 
                   df_historic, amplimix, fl1, fl2, amp2, root, mode, folder)
        #fl1 = 'TEX'
        #fl2 = 'ATTO'
        fl2 = 'TEX'
        fl1 = 'ATTO'
        fig1 = plot_figure_ratio(prefix, df_final, df_ratio, df_uas,
                   df_historic, amplimix, fl1, fl2, amp2, root, mode, folder)
        amp2 = False
        top = Cm(4.64)
        left1 = Cm(0)
        left2 = Cm(16.26)
        height = Cm(12.19)
        blank_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_slide_layout)
        title = slide.shapes.title
        #add rs
        title.text = 'Amplimix: {}'.format(int(amplimix))
        slide.shapes.add_picture(fig1, left1, top, height=height)
        slide.shapes.add_picture(fig2, left2, top, height=height)
###############################################################################
def plot_figure_ratio(prefix, df_final, df_ratio, df_uas, df_historic, amplimix, fl1, fl2, amp2, root, mode, folder):
    '''plot scatter plot figures
    '''
    if amp2:
        uas1 = df_uas.loc[((df_uas['amplimix'].astype(int) == amplimix) &
                           (df_uas['fluorocromo']==fl1) &
                           (df_uas['rs'] == 'am2')), 'uas'].values[0]
        uas2 = df_uas.loc[((df_uas['amplimix'].astype(int) == amplimix) &
                           (df_uas['fluorocromo']==fl2) &
                           (df_uas['rs'] == 'am2')), 'uas'].values[0]
        uas1cp = df_uas.loc[((df_uas['amplimix'].astype(int) == amplimix) &
                           (df_uas['fluorocromo']==fl1) &
                           (df_uas['rs'] == 'am2-cp')), 'uas'].values[0]
        uas2cp = df_uas.loc[((df_uas['amplimix'].astype(int) == amplimix) &
                           (df_uas['fluorocromo']==fl2) &
                           (df_uas['rs'] == 'am2-cp')), 'uas'].values[0]
    else:
        uas1 = df_uas.loc[((df_uas['amplimix'].astype(int) == amplimix) &
                           (df_uas['fluorocromo']==fl1)), 'uas'].values[0]
        uas2 = df_uas.loc[((df_uas['amplimix'].astype(int) == amplimix) &
                           (df_uas['fluorocromo']==fl2)), 'uas'].values[0]
    PREFIX = prefix + '_'
    df_tmp = df_final.loc[(df_final['Amplimix'] == amplimix), :].copy()
    df_tmp_histo = df_historic.loc[(df_historic['Amplimix'] == amplimix), :]
    _, ax = pyplot.subplots()
    col1 = PREFIX + fl1
    col2 = PREFIX + fl2
    df_tmp.loc[df_tmp[col1]<0, col1] = -5
    df_tmp.loc[df_tmp[col2]<0, col2] = -5
    x_min = -10
    y_min = -10
    ax.scatter(df_tmp_histo[col1], df_tmp_histo[col2], color='red')
    ax.scatter(df_tmp[col1], df_tmp[col2])
    ax.plot([x_min, 10000000], [uas2, uas2], color = 'orange', linestyle = ':', linewidth = 1)
    ax.plot([uas1, uas1], [y_min, 10000000], color = 'orange', linestyle = ':', linewidth = 1)
    if amp2:
        ax.plot([x_min, 10000000], [uas2cp, uas2cp], color = 'blue', linestyle = ':', linewidth = 1)
        ax.plot([uas1cp, uas1cp], [y_min, 10000000], color = 'blue', linestyle = ':', linewidth = 1)
    for i, row in df_tmp.iterrows():
        ax.annotate('{}_{}'.format(row['Well'], row['Sample Name']),(row[col1], row[col2]), fontsize=6)
    fig_name = '{}/sc_amp{}_{}_{}.{}.png'.format(folder, amplimix, fl1, fl2, root)
    #plot ratios
    #ratio = '{}/{}'.format(fl1, fl2)
    ratio = '{}/{}'.format(fl2, fl1)
    #corregir això abi
    df_ratio_tmp = df_ratio.loc[((df_ratio['equipo'] == mode) & 
                  (df_ratio['ratio'] == ratio) &
                  (df_ratio['amplimix'].astype(int) == amplimix))].copy()
    min_v = 1
    max_v = 100000000
    step = 1000000
    rs = ''
    if not amp2:
        rs = df_ratio_tmp['rs'].values[0]
        for r in ('down', 'downg', 'up', 'upg'):
            slope = df_ratio_tmp[r].values[0].astype(float)
            #slope = 1/df_ratio_tmp[r].values[0].astype(float)
            x = np.linspace(min_v, max_v, step)
            y = x * slope
            ax.plot(x, y, color='grey')
    else:
        if fl1 == 'FAM':
            rs = 'rs118203905'
        else:
            rs = 'rs118203906'
    histo_xmax = df_tmp_histo[col1].max()
    if histo_xmax is np.nan:
        histo_xmax = 1000
    if histo_xmax == np.inf:
        histo_xmax = 1000
    xmax_ori = df_tmp[col1].max()
    if xmax_ori is np.nan:
        xmax_ori = 1000
    if xmax_ori == np.inf:
        xmax_ori = 1000
    xmax = 1.02 * histo_xmax
    if 1.02 * xmax_ori > xmax:
        xmax = 1.02 * xmax_ori
    if 1.02 * uas1 > xmax:
        xmax = 1.02 * uas1
    histo_ymax = df_tmp_histo[col2].max()
    if histo_ymax is np.nan:
        histo_ymax = 1000
    if histo_ymax == np.inf:
        histo_ymax = 1000
    ymax = 1.02 * histo_ymax
    ymax_ori = df_tmp[col2].max()
    #manage nan and inf values in df_tmp
    if ymax_ori is np.nan:
        ymax_ori = 1000
    if ymax_ori == np.inf:
        ymax_ori = 1000
    if 1.02 * ymax_ori > ymax:
        ymax = 1.02 * ymax_ori
    if 1.02 * uas2 > ymax:
        ymax = 1.02 * uas2
    if ymax > xmax:#arrange a square plot
        xmax = ymax
    else:
        ymax = xmax
    pyplot.ticklabel_format(axis='both', style='sci', scilimits=(0,0))
    #print(ymax, xmax, type(ymax), type(xmax))
    #plot.xlim(0.1, xmax)
    #pyplot.ylim(0.1, ymax)
    pyplot.xlim(10, xmax)
    pyplot.ylim(10, ymax)
    pyplot.xlabel(col1)
    pyplot.ylabel(col2)
    #pyplot.xscale('log')
    #pyplot.yscale('log')
    #pyplot.title('Ratio {}/{} for amplimix{}\nrs: {}'.format(fl1, fl2, amplimix, rs))
    pyplot.title('Ratio {}/{} for amplimix{}\nrs: {}'.format(fl2, fl1, amplimix, rs))
    pyplot.savefig(fig_name)
    pyplot.close()
    return fig_name
###############################################################################
def plot_fluor(prs, prefix, df_final, df_ratio, df_uas, root, mode, folder):
    '''add slides for fluorescence for fluorophore and amplimix
    '''
    fluor = ['FAM', 'HEX', 'TEX', 'ATTO']
    fig_list = []
    for fl in fluor:
        fig = plot_figure_fluor_center(prefix, fl, df_final, root, mode, folder)
        fig_list.append(fig)

    top = Cm(4.14)
    left1 = Cm(0)
    left2 = Cm(16.26)
    height = Cm(14.73)
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    title = slide.shapes.title
    title.text = 'Fluor: {} and {}\n {}:{}'.format(fluor[0], fluor[1], mode, root)
    slide.shapes.add_picture(fig_list[0], left1, top, height=height)
    slide.shapes.add_picture(fig_list[1], left2, top, height=height)

    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    title = slide.shapes.title
    title.text = 'Fluor: {} and {}\n {}:{}'.format(fluor[2], fluor[3], mode, root)
    slide.shapes.add_picture(fig_list[2], left1, top, height=height)
    slide.shapes.add_picture(fig_list[3], left2, top, height=height)
###############################################################################
def plot_figure_fluor_center(prefix, fl, df_tmp, root, mode, folder):
    '''plot final flouroesce for each fluor and for each amplimix
    '''
    PREFIX = prefix + '_'
    fig = pyplot.figure(figsize=(6.4, 5.8))
    _, ax = pyplot.subplots(figsize=(6.4, 5.8))#if you do not change size in here the final png file doesn't have the desired size
    col1 = PREFIX + fl
    df_tmp.loc[df_tmp[col1]<1, col1] = 1
    #plot data
    color_labels = df_tmp['runcode'].unique()
    # List of RGB triplets
    rgb_values = sns.color_palette("Set2", len(color_labels))
    # Map label to RGB
    color_map = dict(zip(color_labels, rgb_values))
    # Finally use the mapped values
    ax.scatter(df_tmp['Amplimix'], df_tmp[col1], c=df_tmp['runcode'].map(color_map))
    #add c+ in x 
    df_tmp_ctrl_pos = df_tmp.loc[df_tmp['Sample Name'] == 'Ctrl +']
    ax.scatter(df_tmp_ctrl_pos['Amplimix'], df_tmp_ctrl_pos[col1], marker= '+', c='blue', label='Ctrl +')
    df_tmp_ctrl_neg = df_tmp.loc[df_tmp['Sample Name'] == 'Ctrl -']
    ax.scatter(df_tmp_ctrl_neg['Amplimix'], df_tmp_ctrl_neg[col1], marker= 'x', c='black', label='Ctrl -')
    # add a legend
    handles = [Line2D([0], [0], marker='o', color='w', markerfacecolor=v, label=k, markersize=8) for k, v in color_map.items()]
    ax.legend(handles=handles, bbox_to_anchor=(1,0), loc="lower right", 
              bbox_transform=fig.transFigure, ncol=5, prop={'size': 6})
    pyplot.subplots_adjust(bottom=0.15)
    if len(color_labels)>30:
        pyplot.subplots_adjust(bottom=0.3)
    elif len(color_labels)>20:
        pyplot.subplots_adjust(bottom=0.25)

    fig_name = '{}/fluor_comp_{}_.{}.png'.format(folder, fl, root)
    ymax = df_tmp[col1].max()
    pyplot.ticklabel_format(axis='both', style='sci', scilimits=(0,0))
    #pyplot.ylim(0.1, ymax)
    pyplot.xlabel('Amplimix      (Ctrl-: \'x\' Ctrl+: \'+\')')
    pyplot.ylabel(col1)
    ax.set_yscale('log')
    pyplot.title('Fluorofor response {} for all amplimixes\n{} @ {}'.format(fl, mode, root))
    pyplot.savefig(fig_name)
    pyplot.close('all')
    return fig_name
###############################################################################
def manage_work(prefix, df_work, df_uas, mode):
    '''work out ratio columns, typing and calling columns
    '''
    #add columns to df_work
    for fluor in ('FAM', 'HEX', 'TEX', 'ATTO'):
        code = fluor + '_uas'
        df_work[code] = np.nan
        code = fluor + '_supera_uas'
        df_work[code] = np.nan
    df_work['NoCall_FAM/HEX'] = np.nan
    df_work['NoCall_TEX/ATTO'] = np.nan
    df_work['Empty'] = np.nan
    order = ['Well', 'Sample Name', 'Amplimix', 
             prefix + '_FAM', 'FAM_uas', 'FAM_supera_uas',
             prefix + '_HEX', 'HEX_uas', 'HEX_supera_uas',  
             'Ratio_FAM/HEX', 'NoCall_FAM/HEX',
             prefix + '_TEX', 'TEX_uas', 'TEX_supera_uas',
             prefix + '_ATTO', 'ATTO_uas', 'ATTO_supera_uas',
             'Ratio_TEX/ATTO', 'NoCall_TEX/ATTO', 'Empty']
    df_work = df_work[order]
    #load columns
    for amplimix in df_uas['amplimix'].unique():
        for fluor in df_uas['fluorocromo'].unique():
            #assign uas to each amplimix
            uas = df_uas.loc[((df_uas['amplimix'] == amplimix) &
                              (df_uas['fluorocromo'] == fluor) &
                              (df_uas['equipo'] == mode )), 'uas'].values
            column = fluor + '_uas'
            df_work.loc[(df_work['Amplimix'].astype(int) == int(amplimix)), column] = uas[0]
            column0 = prefix + '_' + fluor
            column1 = fluor + '_supera_uas'
            df_work.loc[(df_work[column0] > df_work[column]), column1] = 1
            df_work.loc[(df_work[column0] < df_work[column]), column1] = 0
            if amplimix == "2":#correcting pass uas for amp2 ctrl+
                #print(fluor, uas[1])
                #print(df_work.loc[(df_work['Amplimix'] == 2) & (df_work[column0]>uas[1]) & (df_work[column1] == 0)])
                df_work.loc[(df_work['Amplimix'] == 2) & (df_work[column0]>uas[1]) & (df_work[column1] == 0), column1] = 1
                #print(df_work.loc[(df_work['Amplimix'] == 2) & (df_work[column0]>uas[1])])
        column = 'NoCall_FAM/HEX'
        column0 = 'FAM_supera_uas'
        column1 = 'HEX_supera_uas'
        df_work[column] = df_work[column0] + df_work[column1]
        column = 'NoCall_TEX/ATTO'
        column0 = 'TEX_supera_uas'
        column1 = 'ATTO_supera_uas'
        df_work[column] = df_work[column0] + df_work[column1]
        df_work['Empty'] = df_work['NoCall_FAM/HEX'] + df_work['NoCall_TEX/ATTO']
    return df_work
###############################################################################
def load_ratio_from_db(dbfile, mode):
    '''load ratios from db
    '''
    conn = sqlite3.connect(dbfile)
    #load ratios and clean df
    sql = "Select * from thresholds_ratio"
    df_ratio = pd.read_sql(sql, con=conn)
    conn.close()
    #clean df ratio
    df_ratio['amplimix'] = df_ratio['amplimix'].str.replace(r'Amplimix', '')
    df_ratio['ratio'] = df_ratio['ratio'].str.replace(r'TR', 'TEX')
    df_ratio['ratio'] = df_ratio['ratio'].str.replace(r'Cy5', 'ATTO')
    if mode == 'ABI':
        df_ratio['equipo'] = df_ratio['equipo'].str.replace(r'.*' + mode + 
                                                            '.*', mode, regex=True)
        # drop these row indexes from dataFrame 
        index_drop = df_ratio[ df_ratio['equipo'] != mode ].index 
        df_ratio.drop(index_drop, inplace = True) 
    elif mode == 'CFX':
        df_ratio = df_ratio.drop(['tn', 't1', 't2'], axis=1)
    return df_ratio
###############################################################################
def load_uas_from_db(dbfile, mode):
    '''load uas from db
    '''
    conn = sqlite3.connect(dbfile)
    sql = "Select * from thresholds_uas"
    df_uas = pd.read_sql(sql, con=conn)
    conn.close
    #clean df uas
    df_uas['amplimix'] = df_uas['amplimix'].str.replace(r'Amplimix', '')
    df_uas['fluorocromo'] = df_uas['fluorocromo'].str.replace(r'TxR', 'TEX')
    df_uas['fluorocromo'] = df_uas['fluorocromo'].str.replace(r'Cy5', 'ATTO')
    if mode == 'ABI':
        df_uas['equipo'] = df_uas['equipo'].str.replace(r'ABI.*', 'ABI', regex=True)
        index_drop = df_uas[ df_uas['equipo'] != mode ].index 
        # drop these row indexes from dataFrame 
        df_uas.drop(index_drop, inplace = True) 
    elif mode == 'CFX':
        df_uas['equipo'] = mode
        df_uas['version'] = 4
    order = ['version', 'equipo', 'amplimix', 'rs', 'fluorocromo', 'flcode', 'uas']
    df_uas = df_uas[order]
    return df_uas
###############################################################################
def plot_data_plate(prs, root, timestr, df_raw, df_uas, mode, folder):
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = 'curve analysis qpcr'
    content = slide.placeholders[1]
    content.text = 'Input Raw Data:\n {}\nDate: {}'.format(root, timestr)
    plot_sigmoid_per_plate(df_raw,df_uas, prs, mode, folder)
###############################################################################
def plot_sigmoid_per_plate(df_raw, df_uas, prs, mode, folder):
    '''plot all raw data per amplimix and channel
    '''
    amp_list = df_raw['Amplimix'].unique()
    rep_list = df_raw['Reporter'].unique()
    delta = 'ΔRn'
    raw = 'Rn'
    adjust_delta = 'Adjusted Delta'
    dict_adjust_delta = defaultdict(dict)
    dict_delta = defaultdict(dict)
    dict_raw = defaultdict(dict)
    for amp in amp_list:
        for rep in rep_list:
            uas = df_uas.loc[((df_uas['amplimix'].astype(int) == int(amp)) &
                             (df_uas['fluorocromo'] == rep) &
                             (df_uas['rs'].str.contains('cp') == False)), 
                             'uas'].values[0]
            uas2 = 0
            if amp == 2:
                uas2 = df_uas.loc[((df_uas['amplimix'].astype(int) == int(amp)) &
                             (df_uas['fluorocromo'] == rep) &
                             (df_uas['rs'].str.contains('cp') == True)), 
                             'uas'].values[0]
            df_tmp = df_raw.loc[(df_raw['Amplimix'] == amp) & 
                                (df_raw['Reporter'] == rep)]
            dict_adjust_delta[amp][rep] = plot_sigmoid(df_tmp,
                folder, amp, rep, adjust_delta, uas, uas2)
            dict_delta[amp][rep] = plot_sigmoid(df_tmp, folder, amp, rep, delta, uas, uas2)
            dict_raw[amp][rep] = plot_sigmoid(df_tmp, folder, amp, rep, raw, uas, uas2)
    #skip log adjusted delta
    #title = 'Log Adjusted Delta (ΔRn)'
    #add_sigmoid_pptx(prs, dict_adjust_delta, title)
    title = 'Delta (ΔRn)'
    add_sigmoid_pptx(prs, dict_delta, title)
    if mode == 'ABI':
        title = 'Raw Fluorescence Rn'
        add_sigmoid_pptx(prs, dict_raw, title)
    #rmtree(folder)
###############################################################################
def add_sigmoid_pptx(prs, image_dict, title):
    '''add slides for each dict of images
    '''
    rs_dict ={
        1.0 : {'FAM/HEX' : 'rs1799963', 'TEX/ATTO' : 'rs6025'},
        2.0 : {'FAM/HEX' : 'rs118203905', 'TEX/ATTO' : 'rs118203906'},#afegir els que són
        #2.0 : {'FAM/HEX' : 'ampli2.1', 'TEX/ATTO' : 'ampli2.2'},#afegir els que són
        3.0 : {'FAM/HEX' : 'rs1801020', 'TEX/ATTO' : 'rs5985'},
        4.0 : {'FAM/HEX' : 'rs7853989', 'TEX/ATTO' : 'rs8176719'},
        5.0 : {'FAM/HEX' : 'rs8176743', 'TEX/ATTO' : 'rs8176750'},
        6.0 : {'FAM/HEX' : 'rs2232698', 'TEX/ATTO' : 'rs121909548'} 
    }
    #add title slide
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    content = slide.placeholders[1]
    content.text = 'Analysis type:\n{}'.format(title)
    for amp in sorted(image_dict.keys()):
        rep1 = 'FAM'
        rep2 = 'HEX'
        ratio = '{}/{}'.format(rep1, rep2)
        rs = rs_dict[amp][ratio]
        manage_slide_2_fig(prs, amp, rep1, rep2, rs, image_dict[amp][rep1], image_dict[amp][rep2])
        rep1 = 'TEX'
        rep2 = 'ATTO'
        ratio = '{}/{}'.format(rep1, rep2)
        rs = rs_dict[amp][ratio]
        manage_slide_2_fig(prs, amp, rep1, rep2, rs, image_dict[amp][rep1], image_dict[amp][rep2])
###############################################################################
def manage_slide_2_fig(prs, amp, rep1, rep2, rs, img1, img2):
    top = Cm(4.64)
    left1 = Cm(0)
    left2 = Cm(16.26)
    height = Cm(12.19)
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)
    title = slide.shapes.title
    #afegir rs
    title.text = 'Amplimix: {}, rs:{}\n{}\t\t\t\t{}'.format(int(amp), rs, rep1, rep2)
    slide.shapes.add_picture(img1, left1, top, height=height)
    slide.shapes.add_picture(img2, left2, top, height=height)
###############################################################################
def well_to_samp(well):
    num = (np.floor((int(re.sub(r'[a-zA-Z]', '', well))-1) / 6) * 8)
    lett = int(ord(re.sub(r'[0-9]', '', well))) - 64
    ret = num + lett
    return ret
###############################################################################
def sort_well(well):
    num = (np.floor((int(re.sub(r'[a-zA-Z]', '', well[0]))-1) / 6) * 8)
    lett = int(ord(re.sub(r'[0-9]', '', well[0]))) - 64
    ret = num + lett
    return ret
###############################################################################
def plot_sigmoid(df_tmp, folder, amp, rep, column, uas, uas2):
    '''plot sigmoid curve in diff ways
    '''
    _, ax = pyplot.subplots()
    main_sent = ''
    if column == 'Adjusted Delta':
        main_sent = 'log_delta'
    elif column == 'Rn':
        #main_sent = 'log_raw'
        main_sent = 'raw'
    else:
        main_sent = 'delta'
    code = '{}_amplimix_{}_rep_{}'.format(main_sent, int(amp), rep)
    title = '{} amplimix:{} reporter: {}'.format(main_sent, int(amp), rep)
    delta_col = 'Adjusted Delta'
    delta = 'ΔRn'
    # multiline plot with group by
    for well, grp in sorted(df_tmp.groupby(['Well']), key=sort_well):
        #print(well, grp)
        samp = grp['Sample Name'].unique()[0]
        grp[delta_col] = grp[delta]
        grp.loc[grp[delta_col]<1000,'Adjusted Delta']=1000
        lin = 'solid'
        if re.search('Ctrl', samp):
            lin = 'dashdot'
        elif well_to_samp(well) > 8:
            lin = 'dotted'
        code_samp = '{} / {}'.format(samp, well)
        ax.plot(grp['Cycle'], grp[column], label = "{}".format(code_samp), linestyle=lin)
    ax.plot([0, 45], [uas, uas], color='orange', linestyle=':', linewidth=1)
    if amp == 2:
        ax.plot([0, 45], [uas2, uas2], color='blue', linestyle=':', linewidth=1)
    #plot basal 
    ax.axvline(x=3, color='grey', linestyle=':', linewidth=0.5)
    ax.axvline(x=17, color='grey', linestyle=':', linewidth=0.5)
    #location legend
    #https://stackoverflow.com/questions/4700614/how-to-put-the-legend-out-of-the-plot/43439132#43439132
    pyplot.legend(title='sample / well',bbox_to_anchor=(1.02, 1), fontsize='x-small') 
    pyplot.subplots_adjust(right=0.75)
    pyplot.title(title)
    pyplot.xlabel('cycle')
    pyplot.ylabel(main_sent)
    major_ticks = np.arange(4, 45, 5)
    major_ticks = np.insert(major_ticks, 0, 0)
    minor_ticks = np.arange(0, 45, 1)
    ax.set_xticks(major_ticks)
    ax.set_xticks(minor_ticks, minor=True)
    if main_sent == 'log_delta' or main_sent == 'log_raw':
        pyplot.yscale('log')
    fig_name = folder + code + '.png'
    pyplot.savefig(fig_name)
    pyplot.close()
    return fig_name
###############################################################################
def load_raw_ref_from_db(dbfile):
    '''load raw reference from db
    '''
    conn = sqlite3.connect(dbfile)
    #load ratios and clean df
    sql = "Select * from abi_intensidades_ref"
    df_rawref = pd.read_sql(sql, con=conn)
    conn.close()
    #clean df ratio
    df_rawref['amplimix'] = df_rawref['amplimix'].str.replace(r'Amplimix', '')
    df_rawref['flcode'] = df_rawref['flcode'].str.replace(r'0', 'FAM')
    df_rawref['flcode'] = df_rawref['flcode'].str.replace(r'1', 'HEX')
    df_rawref['flcode'] = df_rawref['flcode'].str.replace(r'2', 'TEX')
    df_rawref['flcode'] = df_rawref['flcode'].str.replace(r'3', 'ATTO')
    df_rawref.columns = ['Amplimix', 'Reporter', 'Ref']
    # drop these row indexes from dataFrame 
    return df_rawref
###############################################################################
def slugify(value):
    """
    Normalizes string, converts to lowercase, removes non-alpha characters,
    and converts spaces to hyphens.
    """
    value = unicodedata.normalize('NFKD', value).encode('ascii', 'ignore')
    value = re.sub(r'[^\w\s-]', '', value.decode('utf-8')).strip().lower()
    value = re.sub(r'[-\s]+', '-', value)
    # ...
    return value
###############################################################################
def extract_cfx(prefix, filename):
    '''read, extract and format data from cfx
    '''
    #no data from samples we have to build
    #load data from raw flour
    sheet_2_rep = {0:'FAM', 1:'HEX', 2:'TEX', 3:'ATTO'} 
    df_list = []
    for sheet, rep in sheet_2_rep.items():
        df_tmp = pd.read_excel(filename, header=0, sheet_name=sheet)
        df_tmp = df_tmp.drop([df_tmp.columns[0]], axis=1)#drop column A bc is empty
        df_tmp['Cycle'] = 'c' + df_tmp['Cycle'].astype(str)#change cicle codification add c
        df_tmp = df_tmp.melt(id_vars=['Cycle'], var_name='Well', value_name=prefix)
        #add columns
        df_tmp['Rn'] = df_tmp[prefix] #just to keep identity with abi
        df_tmp['Reporter'] = rep
        df_tmp['Factor'] = 1
        df_tmp['Amplimix'] = df_tmp['Well'].str.replace(r'[a-zA-Z]+', '', regex=True).astype(int)
        #df_tmp.loc[df_tmp['Amplimix'] > 6, 'Amplimix'] = df_tmp['Amplimix'] - 6
        #the order of amplimix is reversed in 2nd half of plate
        df_tmp.loc[df_tmp['Amplimix'] > 6, 'Amplimix'] = 6 - (df_tmp['Amplimix'] - 7)
        #some assays have no sample name, we need to add samples names
        df_tmp['Sample Name'] = np.nan
        df_tmp['Sample Name'] = df_tmp['Sample Name'].mask(pd.isnull, 'code_' +
                                       (((np.floor((df_tmp['Well'].str.replace(
                                         r'[a-zA-Z]+', '', 
                                         regex=True).astype(int) - 1) / 6)) * 8)
                                        #substract 1 to well number -> 0-11
                                        #divide by 6 then apply floor we have 0 o 1
                                        #multiply by 8 --> 0 or 8
                                        + (df_tmp['Well'].str.replace(
                                        r'[0-9]+', '', regex=True).apply(
                                        lambda x:(ord(x)-64)))
                                        - 2).astype(int).astype(str))
                                       #we have well letter to number starts with A -->1
                                       #H --> is 8
                                       #sum previous number we then have 1-16 samples in the same order as plugin does
        df_tmp.loc[(df_tmp['Sample Name'] == 'code_-1'), 'Sample Name'] = 'Ctrl -'
        df_tmp.loc[(df_tmp['Sample Name'] == 'code_0'), 'Sample Name'] = 'Ctrl +'
        order = ['Well', 'Cycle', 'Rn', 'ΔRn', 'Sample Name', 'Reporter', 'Amplimix', 'Factor']
        df_tmp = df_tmp[order]
        df_list.append(df_tmp)
    df_raw = pd.concat(df_list, axis=0, ignore_index=True)
    #Well Sample Name  Amplimix
    df_info = df_raw.drop(df_raw.columns[[1, 2, 3, 5, 7]], axis=1)
    df_info = df_info.drop_duplicates()
    df_filter = df_raw.loc[df_raw['Cycle']=='c45', ('Well', 'Reporter', 'ΔRn')]
    df_filter.loc[(df_filter['ΔRn'] < 1), 'ΔRn'] = 1 # values under 1 set to 1
    df_c45 = df_filter.groupby('Well')[[prefix]].apply(lambda subset: subset.reset_index(drop=True)).unstack().reset_index()
    df_c45 = df_c45.droplevel(level=0, axis=1)
    header = ['Well', prefix + '_FAM', prefix + '_HEX', prefix + '_TEX', prefix + '_ATTO']
    df_c45.columns = header
    df_work = df_info.set_index('Well').join(df_c45.set_index('Well'))
    df_work = df_work.reset_index()
    df_work['Ratio_FAM/HEX'] = df_work[prefix + '_' + 'FAM'] / df_work[prefix + '_' + 'HEX']
    df_work['Ratio_TEX/ATTO'] = df_work[prefix + '_' + 'TEX'] / df_work[prefix + '_' + 'ATTO']
    #['Well', 'Sample Name', 'Amplimix', 'ΔRn_FAM', 'ΔRn_HEX', 'ΔRn_TEX',
    #'ΔRn_ATTO', 'Ratio_FAM/HEX', 'Ratio_TEX/ATTO']
    return df_work, df_raw
###############################################################################
def extract_abi(prefix, filein, dbfile):
    target_name = 'Target Name'
    sample_name = 'Sample Name'
    #load data from samples 
    df_samp = pd.read_excel(filein, index_col=None, header=0,
                            sheet_name=0, skiprows=7)#, keep_default_na=False)
    #load data from raw flour
    df_raw = pd.read_excel(filein, index_col=None, header=0, sheet_name=4, skiprows=7)
    #clean df_samp
    cols = [2,3,4,6,7,9,10,11]
    df_samp.drop(df_samp.columns[cols], axis=1, inplace=True)
    df_tmp = df_samp[[target_name, 'Reporter']].drop_duplicates().sort_values(by=target_name)
    header, reporter_dict = prepare_header(df_tmp, prefix)
    df_samp['Amplimix'] = df_samp['Well'].str.replace(r'[a-zA-Z]+', '', regex=True).astype(int)
    df_samp.loc[df_samp['Amplimix'] > 6, 'Amplimix'] = df_samp['Amplimix'] - 6
    df_raw[sample_name] = np.nan
    df_raw['Reporter'] = np.nan
    df_raw['Amplimix'] = np.nan
    #add by default snumber and well number
    #some assays have no sample name, we need to add samples names
    df_samp['Sample Name'] = df_samp['Sample Name'].mask(pd.isna, 'code_' +
                                       (((np.floor((df_samp['Well'].str.replace(
                                         r'[a-zA-Z]+', '', 
                                         regex=True).astype(int) - 1) / 6)) * 8)
                                        #substract 1 to well number -> 0-11
                                        #divide by 6 then apply floor we have 0 o 1
                                        #multiply by 8 --> 0 or 8
                                     + (df_samp['Well'].str.replace(
                                        r'[0-9]+', '', regex=True).apply(
                                        lambda x:(ord(x)-64)))
                                        -2 ).astype(int).astype(str))
                                       #we have well letter to number starts with A -->1
                                       #H --> is 8
                                       #sum previous number we then have 1-16 samples in the same order as plugin does
    df_samp.loc[(df_samp['Sample Name'] == 'code_-1'), 'Sample Name'] = 'Ctrl -'
    df_samp.loc[(df_samp['Sample Name'] == 'code_0'), 'Sample Name'] = 'Ctrl +'
    for well in df_samp['Well'].unique():
        for target in df_samp[df_samp['Well']== well][target_name].unique():
            sample = df_samp.loc[(df_samp['Well'] == well) & (df_samp[target_name] == target)][sample_name].values
            reporter = df_samp.loc[(df_samp['Well'] == well) & (df_samp[target_name] == target)]['Reporter'].values
            amplimix = df_samp.loc[(df_samp['Well'] == well) & (df_samp[target_name] == target)]['Amplimix'].values
            df_raw.loc[(df_raw['Well'] == well) & 
                       (df_raw[target_name] == target), sample_name] = sample[0]
            df_raw.loc[(df_raw['Well'] == well) & 
                       (df_raw[target_name] == target), 'Reporter'] = reporter_dict[reporter[0]]
            df_raw.loc[(df_raw['Well'] == well) & 
                       (df_raw[target_name] == target), 'Amplimix'] = amplimix[0]
    df_mean = get_factors(df_raw, dbfile)
    ##correct by factor
    df_raw = correct_raw_by_factor(df_raw, df_mean)
    df_samp['c45'] = np.nan
    #df_samp.to_csv('test.csv')
    #load c45
    for index, row in df_samp.iterrows():
        max_cicle = df_raw.loc[((df_raw['Well'] == row['Well']) &
                                (df_raw[target_name] == row[target_name])),
                                'Cycle'].astype(int).max()
        tmp = df_raw.loc[((df_raw['Well']==row['Well']) &
                          (df_raw[target_name]==row[target_name]) & 
                          (df_raw['Cycle']==max_cicle)), prefix].values[0]
        #assign non sigmoid
        #if tmp<1000:
        #    tmp = 1000
        df_samp.loc[index, 'c45'] = tmp
    #drop target name column
    cols = [2]
    df_raw.drop(df_raw.columns[cols], axis=1, inplace=True)
    df_info = df_samp.drop(df_samp.columns[[2,3,5]], axis=1)
    df_info = df_info.drop_duplicates()
    df_c45 = df_samp.groupby('Well')[['c45']].apply(lambda subset: subset.reset_index(drop=True)).unstack().reset_index()
    df_c45 = df_c45.droplevel(level=0, axis=1)
    df_c45.columns = header
    df_work = df_info.set_index('Well').join(df_c45.set_index('Well'))
    df_work = df_work.reset_index()
    df_work['Ratio_FAM/HEX'] = df_work[prefix + '_' + 'FAM'] / df_work[prefix + '_' + 'HEX']
    df_work['Ratio_TEX/ATTO'] = df_work[prefix + '_' + 'TEX'] / df_work[prefix + '_' + 'ATTO']
    return df_work, df_raw
###############################################################################
def correct_raw_by_factor(df_raw, df_mean):
    '''correct by factors
    '''
    for rep in df_mean['Reporter'].unique():
        for amp in df_mean['Amplimix'].unique():
            df_raw.loc[((df_raw['Reporter'] == rep) &
                         (df_raw['Amplimix'].astype(int) == int(amp))),
                         'Factor'] = df_mean.loc[((df_mean['Reporter'] == rep) &
                                 (df_mean['Amplimix'].astype(int) == int(amp))),'Factor'].values[0]
    df_raw.rename(columns = {'Rn':'Ori_Rn', 'ΔRn':'Ori_ΔRn'}, inplace = True) 
    df_raw['Rn'] = df_raw['Ori_Rn'] * df_raw['Factor']
    df_raw['ΔRn'] = df_raw['Ori_ΔRn'] * df_raw['Factor']
    df_raw = df_raw.drop(['Ori_Rn', 'Ori_ΔRn'], axis=1)
    df_raw = df_raw[['Well', 'Cycle', 'Target Name',  'Rn', 'ΔRn', 'Sample Name', 'Reporter', 'Amplimix', 'Factor',]]
    #keep factor?
    return df_raw
###############################################################################
def get_factors(df_raw, dbfile):
    '''get factor by divide reference fluor vs mean
    '''
    df_rawref = load_raw_ref_from_db(dbfile)
    #mean of raw flourescence in the cycles 3-18 
    #get the indexes of non 3-18 cicles
    index_lineal = df_raw[(df_raw['Cycle'].astype(int) < 3) | (df_raw['Cycle'].astype(int) > 18)].index
    #drop other cicles
    df_raw_lineal = df_raw.drop(index_lineal)
    #check filtering is correct
    ##print(df_raw_lineal['Cycle'].unique())
    df_mean = df_raw_lineal.groupby(['Well', 'Reporter', 'Amplimix']).agg({'Rn':['mean', np.std]}).reset_index()
    df_mean = df_mean.droplevel(level=0, axis=1)
    df_mean.columns = ['Well', 'Reporter', 'Amplimix', 'Mean_Rn', 'sd_Rn']
    df_mean['Raw_Ref'] = np.nan
    df_mean['c45_ΔRn'] = np.nan
    #assign raw ref
    for rep in df_mean['Reporter'].unique():
        for amp in df_mean['Amplimix'].unique():
            df_mean.loc[((df_mean['Reporter'] == rep) &
                         (df_mean['Amplimix'].astype(int) == int(amp))),
                         'Raw_Ref'] = df_rawref.loc[((df_rawref['Reporter'] == rep) &
                                 (df_rawref['Amplimix'].astype(int) == int(amp))),'Ref'].values[0]
    #assign c45 Delta Rn
    #find which is the max cycle value
    max_max_cicle = df_raw['Cycle'].astype(int).max()
    if max_max_cicle != 45:
        print('Warning max cicle is diff of 45: {}'.format(max_max_cicle))
    for rep in df_mean['Reporter'].unique():
        for well in df_mean['Well'].unique():
            #for some runs no 45 cicles have been reached, we need to think about limits to accept
            max_cicle = df_raw.loc[((df_raw['Reporter'] == rep) &
                                    (df_raw['Well'] == well)), 'Cycle'].astype(int).max()
            
            df_mean.loc[((df_mean['Reporter'] == rep) &
                         (df_mean['Well'] == well)),
                         'c45_ΔRn'] = df_raw.loc[((df_raw['Reporter'] == rep) &
                                    (df_raw['Well'] == well) &
                                    (df_raw['Cycle'].astype(int) == max_cicle)),'ΔRn'].values[0]
    df_mean['Factor'] = df_mean['Raw_Ref']/df_mean['Mean_Rn']
    return df_mean
###############################################################################
def prepare_header(df_in, pref):
    '''prepare header and check reporter names
    '''
    tmp = df_in['Reporter'].to_list()
    dict_reporter = dict()
    prefix = pref + '_'
    reporter_list = [prefix+'FAM', prefix+'HEX', prefix+'TEX', prefix+'ATTO']
    header = []
    for fluor in tmp:
        if re.search('TEX', fluor) or re.search('TXR', fluor):
            dict_reporter[fluor] = 'TEX'
            fluor =  prefix +'TEX'
        elif re.search('VIC', fluor):
            dict_reporter[fluor] = 'HEX'
            fluor = prefix + 'HEX'
        elif re.search('CY5', fluor):
            dict_reporter[fluor] = 'ATTO'
            fluor = prefix + 'ATTO'
        else:
            dict_reporter[fluor] = fluor
            fluor = prefix + fluor
        header.append(fluor)
    count = 0
    correct = 0
    for idx, rep in enumerate (header):
        if rep not in reporter_list:
            count += 1
            correct = idx
            print('Warning: this reporter: {} not expected'.format(rep))
        else:
            reporter_list.remove(rep)
    if count == 1:
        print('Info: correcting this reporter: {} to {}'.format(
              header[correct], reporter_list[0]))
        dict_reporter[re.sub(prefix, '', header[correct])] = re.sub(prefix, '', reporter_list[0])
        header[correct] = reporter_list[0]
    elif count >0:
        print('Error: problems with reporter names')
        print(header)
        sys.exit()
    reheader = ['Well']
    reheader.extend(header)
    return reheader, dict_reporter
###############################################################################
def parse_arguments():
    '''parsing input arguments
    '''
    args = ''
    parser = argparse.ArgumentParser(description="import fuck file to process")
    parser.add_argument("-i", "--infile", type=argparse.FileType('r'),
                        required=False, help="input file in eds format")
    parser.add_argument("-f", "--folder", #type=argparse.FileType('r'),
                        required=False, 
                        help="data folder if any, otherwise pwd")
    parser.add_argument("-t", "--plate_type", default='RUN', const='RUN',
                        nargs = '?', choices=('RUN', 'SETUP'),
                        help='Type of plate, to be stored')

    args = parser.parse_args()
    return(args)
###############################################################################
if __name__ == '__main__':
    ARGUM = parse_arguments()
    main(ARGUM)
